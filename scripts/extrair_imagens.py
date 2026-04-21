"""
Extrai imagens de produto do tabela_DL.pptx e organiza em
src/imagens/produtos/{categoria-slug}/ com nomes baseados no texto do slide.
Uso: py scripts/extrair_imagens.py
"""

import json
import re
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = Path(__file__).parent.parent / "data" / "tabela_DL.pptx"
IMGS_ROOT  = Path(__file__).parent.parent / "src" / "imagens" / "produtos"
JSON_PATH  = Path(__file__).parent.parent / "data" / "produtos.json"

IGNORE_TEXTS = {
    "DL", "FOTO", "DL DESCARTÁVEIS", "CATÁLOGO DE PRODUTOS",
    "PIX", "CNPJ", "C6 BANK", "DL DESCART", "FOTO\ufffd\ufffd",
}

FOOTER_PARTIAL = ["DL DESCART", "Cat\u00e1logo de Produtos", "Janeiro", "Abril"]


def to_slug(text: str) -> str:
    text = unicodedata.normalize("NFD", text)
    text = "".join(c for c in text if unicodedata.category(c) != "Mn")
    text = text.lower().strip()
    text = re.sub(r"[^a-z0-9]+", "-", text)
    return text.strip("-")


def is_footer(text: str) -> bool:
    return any(p in text for p in FOOTER_PARTIAL)


def get_slide_texts(slide) -> list[str]:
    texts = []
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and not is_footer(t) and t.upper() not in IGNORE_TEXTS:
                        texts.append(t)
        except Exception:
            pass
    return texts


def get_category_slug(title: str) -> str:
    """Converte título do slide em slug para casar com produtos.json."""
    return to_slug(title)


def get_product_images(slide) -> list:
    """Retorna todas as shapes do tipo PICTURE no slide."""
    images = []
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append(shape)
        except Exception:
            pass
    return images


def load_products_by_category() -> dict:
    """Carrega produtos.json e indexa por slug de categoria."""
    if not JSON_PATH.exists():
        return {}
    data = json.loads(JSON_PATH.read_text(encoding="utf-8"))
    result = {}
    for cat in data.get("categorias", []):
        result[cat["slug"]] = [p["nome"] for p in cat["produtos"]]
    return result


def deduce_product_names_from_slide(slide, category_slug: str) -> list[str]:
    """
    Tenta deduzir nomes de produtos a partir dos textos do slide,
    filtrando dimensões (ex: '25 x 35 cm') e textos genéricos.
    """
    texts = get_slide_texts(slide)
    dim_pattern = re.compile(r"^\d+\s*[xX×]\s*\d+", )
    products = []
    for t in texts:
        if dim_pattern.match(t):
            continue
        if re.match(r"^c/\d+", t, re.I):
            continue
        # Evitar nomes de categoria (geralmente uppercase longo)
        if t.upper() == t and len(t) > 20 and "(" not in t:
            continue
        products.append(t)
    return products


def main():
    print(f"Abrindo {PPTX_PATH} …")
    prs = Presentation(str(PPTX_PATH))
    products_by_cat = load_products_by_category()

    total_saved = 0
    skipped_slides = 0

    for slide_idx, slide in enumerate(prs.slides):
        texts = get_slide_texts(slide)
        if not texts:
            skipped_slides += 1
            continue

        # Primeiro texto não-genérico é o título da categoria
        category_title = texts[0] if texts else None
        if not category_title or category_title.upper() in IGNORE_TEXTS:
            skipped_slides += 1
            continue

        # Slide de capa (slide 0)
        if slide_idx == 0:
            skipped_slides += 1
            continue

        cat_slug = get_category_slug(category_title)

        # Remover sufixos como " (1)", " (2)" do slug para casar com JSON
        base_slug = re.sub(r"-\d+$", "", cat_slug)

        images = get_product_images(slide)
        if not images:
            skipped_slides += 1
            continue

        # Criar pasta da categoria
        cat_dir = IMGS_ROOT / base_slug
        cat_dir.mkdir(parents=True, exist_ok=True)

        # Deduzir nomes de produtos do slide para nomear arquivos
        product_names = deduce_product_names_from_slide(slide, base_slug)

        print(f"\nSlide {slide_idx + 1}: {category_title} => {base_slug} ({len(images)} imgs)")

        for img_idx, img_shape in enumerate(images):
            try:
                image = img_shape.image
                ext = image.ext  # 'png', 'jpeg', etc.

                # Tentar usar nome do produto correspondente
                if img_idx < len(product_names):
                    name = to_slug(product_names[img_idx])
                else:
                    name = f"produto-{img_idx + 1}"

                filename = f"{name}.{ext}"
                dest = cat_dir / filename

                dest.write_bytes(image.blob)
                total_saved += 1
                print(f"  [{img_idx + 1}] {filename}")
            except Exception as e:
                print(f"  [{img_idx + 1}] ERRO: {e}")

    print(f"\n{total_saved} imagens extraidas | {skipped_slides} slides ignorados")
    print(f"Destino: {IMGS_ROOT}")


if __name__ == "__main__":
    main()
